"""Generate repo images: pipeline diagram (16:9) and social preview (2:1)."""

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "docs" / "images"
OUT.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x14, 0x2E, 0x52)
DEEP = RGBColor(0x0A, 0x1B, 0x33)
DEEP2 = RGBColor(0x10, 0x2A, 0x4E)
BLUE = RGBColor(0x00, 0x78, 0xD4)
BLUE_LT = RGBColor(0x6C, 0xBE, 0xFF)
AMBER = RGBColor(0xD8, 0x7A, 0x00)
GREEN = RGBColor(0x24, 0x8A, 0x52)
GREEN_LT = RGBColor(0x6A, 0xD0, 0x98)
INK = RGBColor(0x1E, 0x28, 0x33)
MUTED = RGBColor(0x5E, 0x69, 0x77)
FAINT = RGBColor(0x9A, 0xAC, 0xC4)
PALE = RGBColor(0xEF, 0xF5, 0xFB)
LINE = RGBColor(0xD3, 0xDF, 0xEC)
RING_DK = RGBColor(0x1C, 0x39, 0x63)
PBI_YELLOW = RGBColor(0xF2, 0xC8, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DISPLAY = "Aptos Display"
BODY = "Aptos"


def _fill(sp, c):
    if c is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = c


def _line(sp, c, w=1.0):
    if c is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = c; sp.line.width = Pt(w)


def box(s, kind, x, y, w, h, f, ln=None, lw=1.0, adj=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(sp, f); _line(sp, ln, lw); sp.shadow.inherit = False
    if adj is not None and sp.adjustments:
        sp.adjustments[0] = adj
    return sp


def grad(sp, stops, ang):
    spPr = sp._element.spPr
    for t in ("a:solidFill", "a:noFill", "a:gradFill"):
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)
    gs = "".join(f'<a:gs pos="{int(p*100000)}"><a:srgbClr val="{c}"><a:alpha val="{int(a*100000)}"/></a:srgbClr></a:gs>' for p, c, a in stops)
    spPr.find(qn("a:prstGeom")).addnext(parse_xml(f'<a:gradFill {nsdecls("a")}><a:gsLst>{gs}</a:gsLst><a:lin ang="{int(ang*60000)}" scaled="1"/></a:gradFill>'))
    _line(sp, None)


def txt(s, x, y, w, h, runs, anchor=MSO_ANCHOR.MIDDLE):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, r in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = r.get("align", PP_ALIGN.LEFT); p.space_after = Pt(0); p.space_before = Pt(0)
        run = p.add_run(); run.text = r["t"]
        run.font.name = r.get("font", DISPLAY); run.font.size = Pt(r["sz"])
        run.font.bold = r.get("b", False); run.font.color.rgb = r["c"]
    return tb


def pbi_icon(s, x, y, h):
    bw, gap = h * 0.24, h * 0.15
    for i, f in enumerate((0.42, 0.62, 0.82, 1.0)):
        box(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + i * (bw + gap), y + (h - h * f), bw, h * f, PBI_YELLOW, None, 1.0, 0.4)


def dark_bg(s, w, h):
    r = box(s, MSO_SHAPE.RECTANGLE, 0, 0, w, h, DEEP)
    grad(r, [(0.0, "0A1B33", 1.0), (0.55, "10294C", 1.0), (1.0, "16385F", 1.0)], 55)
    box(s, MSO_SHAPE.OVAL, w - 3.3, -2.3, 6.6, 6.6, None, RING_DK, 1.25)
    box(s, MSO_SHAPE.OVAL, -3.4, h - 3.4, 6.6, 6.6, None, RING_DK, 1.25)


# ---------------------------------------------------------------- pipeline 16:9
def build_pipeline():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    box(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, BLUE)
    pbi_icon(s, 0.6, 0.5, 0.5)
    txt(s, 1.4, 0.5, 5, 0.5, [{"t": "Power BI", "sz": 20, "c": NAVY, "b": True}])
    txt(s, 0.6, 1.15, 12.1, 0.9, [{"t": "AI-Readiness Pipeline \u2014 6 prompt-driven steps", "sz": 30, "c": NAVY, "b": True}])
    steps = [("1", "CLEANUP", "Remove dead\nobjects", BLUE),
             ("2", "OPTIMIZE", "Star schema\n& DAX", AMBER),
             ("3", "BPA", "AI-readiness\nrules", GREEN),
             ("4", "DESCRIBE", "Descriptions\n+ synonyms", BLUE),
             ("5", "RENAME", "Business names\nhide keys", AMBER),
             ("6", "SCORE", "0\u2013100\nreadiness", GREEN)]
    n = len(steps); gap = 0.35
    cw = (12.1 - gap * (n - 1)) / n
    x0 = 0.6
    for i, (num, name, body, accent) in enumerate(steps):
        x = x0 + i * (cw + gap)
        box(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.7, cw, 2.5, WHITE, LINE, 1.25, 0.08)
        box(s, MSO_SHAPE.OVAL, x + cw / 2 - 0.42, 3.0, 0.84, 0.84, accent)
        txt(s, x + cw / 2 - 0.42, 2.98, 0.84, 0.84, [{"t": num, "sz": 24, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER}])
        txt(s, x, 4.0, cw, 0.5, [{"t": name, "sz": 15, "c": accent, "b": True, "align": PP_ALIGN.CENTER}])
        txt(s, x, 4.5, cw, 0.7, [{"t": body, "sz": 11.5, "c": INK, "align": PP_ALIGN.CENTER, "font": BODY}])
        if i < n - 1:
            box(s, MSO_SHAPE.RIGHT_ARROW, x + cw + 0.02, 3.75, gap - 0.04, 0.35, BLUE)
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 5.7, 12.1, 1.0, PALE, LINE, 1.0, 0.12)
    txt(s, 0.9, 5.7, 11.5, 1.0, [
        {"t": "You prompt  \u00b7  the agent builds (GitHub Copilot + Power BI Modeling MCP + skills)  \u00b7  you approve", "sz": 14, "c": NAVY, "b": True}])
    p = ROOT / "_pipeline.pptx"; prs.save(p); return p


# ---------------------------------------------------------------- social 2:1
def build_social():
    prs = Presentation(); prs.slide_width = Inches(12.8); prs.slide_height = Inches(6.4)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = DEEP
    dark_bg(s, 12.8, 6.4)
    pbi_icon(s, 0.7, 0.7, 0.6)
    txt(s, 1.7, 0.68, 6, 0.6, [{"t": "Power BI", "sz": 24, "c": WHITE, "b": True}])
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 1.6, 4.0, 0.6, PBI_YELLOW, None, 1.0, 0.5)
    txt(s, 0.7, 1.6, 4.0, 0.6, [{"t": "POWER BI + COPILOT", "sz": 14, "c": DEEP, "b": True, "align": PP_ALIGN.CENTER, "font": BODY}])
    txt(s, 0.7, 2.5, 8.5, 2.2, [
        {"t": "Power BI \u2192 AI-Ready", "sz": 44, "c": WHITE, "b": True},
        {"t": "in 6 prompt-driven steps", "sz": 30, "c": BLUE_LT, "b": True}])
    txt(s, 0.72, 5.2, 9.0, 0.7, [
        {"t": "Cleanup \u2192 Optimize \u2192 BPA \u2192 Describe \u2192 Rename \u2192 Score", "sz": 15, "c": FAINT, "font": BODY}])
    box(s, MSO_SHAPE.OVAL, 9.3, 1.9, 2.6, 2.6, None, RING_DK, 5.0)
    orb = box(s, MSO_SHAPE.OVAL, 9.6, 2.2, 2.0, 2.0, DEEP2, BLUE_LT, 1.5)
    grad(orb, [(0.0, "173D6B", 1.0), (1.0, "0C2244", 1.0)], 90)
    txt(s, 9.6, 2.75, 2.0, 0.9, [{"t": "13 \u2192 98", "sz": 26, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER},
                                 {"t": "AI-Readiness", "sz": 11, "c": FAINT, "align": PP_ALIGN.CENTER, "font": BODY}])
    p = ROOT / "_social.pptx"; prs.save(p); return p


print(build_pipeline())
print(build_social())
